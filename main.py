import PySimpleGUI as sg
import win32com.client
import time
import os
from pptx import Presentation
user = os.environ.get('USERNAME')
FOLDER_PATH = fr"C:\Users\{user}\Documents\birthday_videos"
PRESENTATION_PATH = os.path.join(FOLDER_PATH, "bday.pptx")


def window_main():
    sg.theme('DarkAmber')  # Keep things interesting for your users
    layout = [[sg.Text('Good morning sir, who celebrates today?')],
              [sg.Text("Name:"), sg.Input(key='name')],
              [sg.Text("Age:"), sg.Combo(list(range(120)), key='age')],
              [sg.Button('Make Video'), sg.Exit()]]

    window = sg.Window('Birthday Video Generator', layout, size=(400,400))

    while True:  # The Event Loop
        event, values = window.read()
        change_slide(name=values['name'], age=values['age'])
        vid_name = f"{values['name']}_{values['age']}_birthday_vid"
        make_video(vid_name=vid_name)
        sg.popup("the video is ready!")
        if event == sg.WIN_CLOSED or event == 'Exit':
            break


    window.close()


def change_slide(name, age):
    text_dic = {0 : f" בוקר טוב {name}", 2 : f"בשמי ובשם מפקדי וחיילי חטיבת האש 214, ברצוני לאחל לך מזל טוב לרגל יום הולדתך ה{age}!"}
    prs = Presentation(PRESENTATION_PATH)
    # To get shapes in your slides
    slides = [slide for slide in prs.slides]
    slide = slides[0]
    shapes = slide.shapes
    text_index = 0
    for shape in shapes:
      if shape.has_text_frame:
          if text_index != 3 and text_index != 1:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for r_index, run in enumerate(paragraph.runs):
                    if r_index == 0:
                        run.text = text_dic[text_index]
                    else:
                        run.text = ''
          text_index += 1
    prs.save(PRESENTATION_PATH)


def make_video(vid_name: str):
    # Start converting
    ppt = win32com.client.Dispatch('PowerPoint.Application')
    presentation = ppt.Presentations.Open(PRESENTATION_PATH, WithWindow=False)
    temp_path = f"{FOLDER_PATH}\{vid_name}.mp4"
    presentation.SaveAs(temp_path, 39)
    while True:
        if os.path.exists(temp_path):
            try:
                if os.path.getsize(temp_path) > 0:
                    break
            except Exception as e:
                print(e)
    presentation.close()




# Press the green button in the gutter to run the script.
if __name__ == '__main__':
    window_main()


